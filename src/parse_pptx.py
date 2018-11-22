import os
from pptx import Presentation

def parse_pptx(doc_path):


	prs = Presentation(doc_path)

	text_runs = []

	for slide in prs.slides:
	    for shape in slide.shapes:
	        if not shape.has_text_frame:
	            continue
	        for paragraph in shape.text_frame.paragraphs:
	            for run in paragraph.runs:
	                text_runs.append(run.text)

	parsed_text=' '.join(text_runs)

	return (parsed_text)

